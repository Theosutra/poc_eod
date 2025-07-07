#!/usr/bin/env python3
"""
Pipeline RAG principal pour EODEN
Orchestration complète : Drive → Chunking → Embeddings → Vector Store
"""

import os
import sys
import logging
from pathlib import Path
from typing import Dict, List, Any
import yaml
from dotenv import load_dotenv
from tqdm import tqdm

# Imports locaux
from drive_connector import DriveConnector
from simplified_embedding_manager import SimplifiedEmbeddingManager, SimplifiedChunk
from vector_store import VectorManager, VectorDocument

# Configuration du logging
import os
os.makedirs("logs", exist_ok=True)  # Créer le dossier logs s'il n'existe pas

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('logs/poc.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class RAGPipeline:
    """Pipeline principal RAG pour EODEN"""
    
    def __init__(self, config_path: str = "config/settings.yaml"):
        # Charger la configuration
        load_dotenv()
        self.config = self._load_config(config_path)
        
        # Initialiser les composants
        self.drive_connector = None
        self.embedding_manager = None
        self.vector_manager = None
        
        self._setup_directories()
        self._init_components()
    
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Charger la configuration YAML"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = yaml.safe_load(f)
            
            # Substituer les variables d'environnement
            config = self._substitute_env_vars(config)
            logger.info(f"Configuration chargée: {config_path}")
            return config
            
        except Exception as e:
            logger.error(f"Erreur chargement config: {e}")
            sys.exit(1)
    
    def _substitute_env_vars(self, obj):
        """Substituer les variables d'environnement dans la config"""
        if isinstance(obj, dict):
            return {k: self._substitute_env_vars(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [self._substitute_env_vars(item) for item in obj]
        elif isinstance(obj, str) and obj.startswith('${') and obj.endswith('}'):
            env_var = obj[2:-1]
            return os.getenv(env_var, obj)
        else:
            return obj
    
    def _setup_directories(self):
        """Créer les répertoires nécessaires"""
        directories = [
            "data/cache",
            "data/embeddings", 
            "data/output",
            "logs",
            "config"
        ]
        
        for dir_path in directories:
            Path(dir_path).mkdir(parents=True, exist_ok=True)
    
    def _init_components(self):
        """Initialiser les composants du pipeline"""
        try:
            # Google Drive Connector
            logger.info("🔗 Initialisation Google Drive...")
            self.drive_connector = DriveConnector(
                credentials_path=self.config['google_drive']['credentials_file'],
                cache_dir=self.config['google_drive']['cache_directory']
            )
            
            # Simplified Embedding Manager
            logger.info("🧠 Initialisation Simplified Embedding Manager...")
            gemini_api_key = self.config['apis']['gemini_api_key']
            if not gemini_api_key:
                raise ValueError("GEMINI_API_KEY manquant dans .env")
            
            self.embedding_manager = SimplifiedEmbeddingManager(
                api_key=gemini_api_key,
                cache_dir=self.config['embeddings']['cache_directory']
            )
            
            # Vector Store Manager
            logger.info("📦 Initialisation Vector Store...")
            self.vector_manager = VectorManager(config_path="config/settings.yaml")
            
            logger.info("✅ Tous les composants initialisés")
            
        except Exception as e:
            logger.error(f"❌ Erreur initialisation: {e}")
            sys.exit(1)
    
    def discover_documents(self) -> List[Dict[str, Any]]:
        """Découvrir les documents dans Google Drive"""
        logger.info("🔍 Découverte des documents Google Drive...")
        
        all_files = []
        
        # Parcourir les dossiers configurés
        for folder_config in self.config['google_drive']['source_folders']:
            folder_name = folder_config['name']
            recursive = folder_config.get('recursive', True)
            
            logger.info(f"📁 Exploration du dossier: {folder_name}")
            
            # Trouver le dossier
            folder_id = self.drive_connector.get_folder_by_name(folder_name)
            if not folder_id:
                logger.warning(f"⚠️ Dossier non trouvé: {folder_name}")
                continue
            
            # Lister les fichiers
            files = self.drive_connector.list_files(folder_id, recursive=recursive)
            all_files.extend(files)
            logger.info(f"✅ {len(files)} fichiers trouvés dans {folder_name}")
        
        # Filtrer par extensions supportées
        supported_exts = self.config['google_drive']['supported_extensions']
        filtered_files = []
        
        for file_info in all_files:
            file_name = file_info['name']
            if any(file_name.lower().endswith(ext.lower()) for ext in supported_exts):
                filtered_files.append(file_info)
        
        logger.info(f"🎯 {len(filtered_files)} fichiers supportés découverts")
        return filtered_files
    
    def download_documents(self, files_info: List[Dict[str, Any]]) -> Dict[str, str]:
        """Télécharger et extraire le contenu des documents"""
        logger.info("⬇️ Téléchargement des documents...")
        
        downloaded_contents = {}
        
        with tqdm(total=len(files_info), desc="Téléchargement") as pbar:
            for file_info in files_info:
                try:
                    # Télécharger le fichier
                    cache_path = self.drive_connector.download_file(
                        file_info['id'],
                        file_info['name'],
                        file_info['mimeType']
                    )
                    
                    if cache_path:
                        # Extraire le contenu
                        content = self._extract_content(cache_path)
                        if content:
                            downloaded_contents[str(cache_path)] = content
                            logger.debug(f"✅ Contenu extrait: {file_info['name']}")
                        else:
                            logger.warning(f"⚠️ Extraction échouée: {file_info['name']}")
                    
                except Exception as e:
                    logger.error(f"❌ Erreur téléchargement {file_info['name']}: {e}")
                
                pbar.update(1)
        
        logger.info(f"✅ {len(downloaded_contents)} documents téléchargés et extraits")
        return downloaded_contents
    
    def _extract_content(self, file_path: Path) -> str:
        """Extraire le contenu textuel d'un fichier"""
        try:
            file_ext = file_path.suffix.lower()
            
            if file_ext == '.pdf':
                return self._extract_pdf_content(file_path)
            elif file_ext == '.docx':
                return self._extract_docx_content(file_path)
            elif file_ext == '.pptx':
                return self._extract_pptx_content(file_path)
            elif file_ext == '.xlsx':
                return self._extract_xlsx_content(file_path)
            else:
                logger.warning(f"Type de fichier non supporté: {file_ext}")
                return ""
                
        except Exception as e:
            logger.error(f"Erreur extraction {file_path}: {e}")
            return ""
    
    def _extract_pdf_content(self, file_path: Path) -> str:
        """Extraire le contenu d'un PDF avec OCR des images"""
        import pymupdf  # fitz
        
        text = ""
        doc = pymupdf.open(str(file_path))
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            text += f"\n--- Page {page_num + 1} ---\n"
            
            # Extraire le texte normal
            page_text = page.get_text()
            text += page_text
            
            # Si peu de texte, essayer l'OCR sur les images
            if len(page_text.strip()) < 100:
                try:
                    ocr_text = self._extract_images_with_ocr(page)
                    if ocr_text:
                        text += f"\n[OCR] {ocr_text}"
                        logger.info(f"OCR appliqué sur page {page_num + 1}")
                except Exception as e:
                    logger.warning(f"Erreur OCR page {page_num + 1}: {e}")
        
        doc.close()
        return text.strip()
    
    def _extract_docx_content(self, file_path: Path) -> str:
        """Extraire le contenu d'un Word"""
        from docx import Document
        
        doc = Document(str(file_path))
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extraire les tableaux
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                text += row_text + "\n"
        
        return text.strip()
    
    def _extract_pptx_content(self, file_path: Path) -> str:
        """Extraire le contenu d'un PowerPoint avec OCR des images"""
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            text += slide_text
            
            # Si peu de texte sur la slide, essayer l'OCR sur les images
            if len(slide_text.strip()) < 50:
                try:
                    ocr_text = self._extract_slide_images_with_ocr(slide)
                    if ocr_text:
                        text += f"\n[OCR] {ocr_text}"
                        logger.info(f"OCR appliqué sur slide {slide_num}")
                except Exception as e:
                    logger.warning(f"Erreur OCR slide {slide_num}: {e}")
        
        return text.strip()
    
    def _extract_xlsx_content(self, file_path: Path) -> str:
        """Extraire le contenu d'un Excel"""
        import pandas as pd
        
        text = ""
        
        try:
            # Lire toutes les feuilles
            excel_file = pd.ExcelFile(str(file_path))
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(str(file_path), sheet_name=sheet_name)
                
                text += f"\n--- Feuille: {sheet_name} ---\n"
                
                # Convertir en texte structuré
                for index, row in df.iterrows():
                    row_text = " | ".join(str(value) for value in row.values if pd.notna(value))
                    if row_text.strip():
                        text += row_text + "\n"
        
        except Exception as e:
            logger.warning(f"Erreur lecture Excel {file_path}: {e}")
        
        return text.strip()
    
    def _extract_images_with_ocr(self, page) -> str:
        """Extraire le texte des images d'une page PDF avec OCR"""
        try:
            # Méthode 1: Gemini Vision (intelligent, même API)
            return self._ocr_with_gemini_vision(page)
        except Exception as e:
            logger.warning(f"Gemini Vision OCR failed: {e}")
            try:
                # Méthode 2: Tesseract (fallback gratuit)
                return self._ocr_with_tesseract(page)
            except Exception as e2:
                logger.warning(f"Tesseract OCR failed: {e2}")
                try:
                    # Méthode 3: Google Vision API (fallback précis)
                    return self._ocr_with_google_vision(page)
                except Exception as e3:
                    logger.warning(f"Google Vision OCR failed: {e3}")
                    return ""
    
    def _ocr_with_gemini_vision(self, page) -> str:
        """OCR avec Gemini Vision (intelligent et contextuel)"""
        try:
            import google.generativeai as genai
            import base64
            import io
            
            # Extraire les images de la page
            image_list = page.get_images(full=True)
            ocr_texts = []
            
            for img_index, img in enumerate(image_list):
                try:
                    # Récupérer l'image
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Encoder en base64 pour Gemini
                    image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                    
                    # Prompt optimisé pour Gemini 2.0 Flash
                    prompt = """
                    Vous êtes un expert en analyse de documents d'investissement. Analysez cette image avec une précision maximale.
                    
                    MISSIONS:
                    1. OCR complet: Extraire TOUT le texte visible avec une précision parfaite
                    2. Analyse financière: Identifier toutes les métriques, ratios, évolutions
                    3. Structuration intelligente: Organiser les données pour faciliter l'analyse
                    4. Insights business: Dégager les tendances et signaux clés
                    
                    FORMAT DE SORTIE:
                    
                    TEXTE EXTRAIT:
                    [Transcription exacte de tout le texte visible]
                    
                    MÉTRIQUES FINANCIÈRES:
                    [Lister tous les chiffres avec unités: CA, EBITDA, croissance, ratios, etc.]
                    
                    DONNÉES TEMPORELLES:
                    [Évolutions, comparaisons annuelles, tendances identifiées]
                    
                    INSIGHTS CLÉS:
                    [Points saillants pour un investisseur: performance, risques, opportunités]
                    
                    TYPE DE CONTENU:
                    [Graphique/Tableau/Diagramme/Schéma + description technique]
                    
                    Soyez exhaustif et précis. Gemini 2.5 Excellence attendue.
                    """
                    
                    # Créer le modèle Gemini Vision 2.5
                    model = genai.GenerativeModel('gemini-2.5-flash')
                    
                    # Analyser l'image
                    response = model.generate_content([
                        prompt,
                        {
                            "mime_type": f"image/{base_image.get('ext', 'png')}",
                            "data": image_b64
                        }
                    ])
                    
                    if response.text and response.text.strip():
                        analyzed_content = response.text.strip()
                        ocr_texts.append(f"[Image {img_index + 1} - Gemini Analysis]\n{analyzed_content}")
                        logger.info(f"Gemini Vision OCR réussi pour image {img_index + 1}")
                        
                except Exception as e:
                    logger.warning(f"Erreur Gemini Vision image {img_index}: {e}")
                    continue
            
            return "\n\n".join(ocr_texts)
            
        except Exception as e:
            logger.warning(f"Gemini Vision OCR non disponible: {e}")
            raise e
    
    def _ocr_with_tesseract(self, page) -> str:
        """OCR avec Tesseract (gratuit)"""
        try:
            import pytesseract
            from PIL import Image
            import io
            
            # Extraire les images de la page
            image_list = page.get_images(full=True)
            ocr_texts = []
            
            for img_index, img in enumerate(image_list):
                try:
                    # Récupérer l'image
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Convertir en PIL Image
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    # OCR avec Tesseract
                    text = pytesseract.image_to_string(image, lang='fra+eng')
                    
                    if text.strip():
                        ocr_texts.append(f"[Image {img_index + 1}] {text.strip()}")
                        
                except Exception as e:
                    logger.warning(f"Erreur OCR image {img_index}: {e}")
                    continue
            
            return "\n".join(ocr_texts)
            
        except ImportError:
            logger.warning("pytesseract non installé. Installer avec: pip install pytesseract")
            return ""
    
    def _ocr_with_google_vision(self, page) -> str:
        """OCR avec Google Vision API (plus précis mais payant)"""
        try:
            from google.cloud import vision
            import io
            
            client = vision.ImageAnnotatorClient()
            image_list = page.get_images(full=True)
            ocr_texts = []
            
            for img_index, img in enumerate(image_list):
                try:
                    # Récupérer l'image
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Analyser avec Google Vision
                    image = vision.Image(content=image_bytes)
                    response = client.text_detection(image=image)
                    
                    if response.text_annotations:
                        text = response.text_annotations[0].description
                        ocr_texts.append(f"[Image {img_index + 1}] {text.strip()}")
                        
                except Exception as e:
                    logger.warning(f"Erreur Google Vision image {img_index}: {e}")
                    continue
            
            return "\n".join(ocr_texts)
            
        except ImportError:
            logger.warning("Google Cloud Vision non installé. Installer avec: pip install google-cloud-vision")
            return ""
    
    def _extract_slide_images_with_ocr(self, slide) -> str:
        """Extraire le texte des images d'une slide PowerPoint avec Gemini Vision"""
        try:
            import google.generativeai as genai
            import base64
            import io
            from PIL import Image
            
            ocr_texts = []
            
            # Parcourir les formes de la slide
            for shape_index, shape in enumerate(slide.shapes):
                # Vérifier si c'est une image
                if hasattr(shape, 'image'):
                    try:
                        # Méthode 1: Gemini Vision
                        try:
                            # Extraire l'image
                            image_bytes = shape.image.blob
                            image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            
                            # Prompt optimisé pour slides avec Gemini 2.0
                            prompt = """
                            Analysez cette image de slide de présentation business avec la puissance de Gemini 2.0.
                            
                            EXTRACTION COMPLÈTE:
                            1. OCR exhaustif de tout le texte (titres, sous-titres, légendes, annotations)
                            2. Données quantitatives: chiffres, pourcentages, ratios, dates
                            3. Éléments visuels: graphiques, schémas, flèches, codes couleur
                            4. Message stratégique: objectif de la slide pour l'audience
                            
                            STRUCTURE DE RÉPONSE:
                            
                            CONTENU TEXTUEL:
                            [Transcription complète et structurée]
                            
                            DONNÉES BUSINESS:
                            [Métriques, KPIs, projections, comparatifs]
                            
                            ÉLÉMENTS VISUELS:
                            [Description des graphiques, schémas, diagrammes]
                            
                            MESSAGE STRATÉGIQUE:
                            [Objectif de communication, points clés pour investisseurs]
                            
                            CONTEXTE SLIDE:
                            [Position probable dans la présentation, audience cible]
                            
                            Exploitez pleinement Gemini 2.5 pour une analyse approfondie.
                            """
                            
                            model = genai.GenerativeModel('gemini-2.5-flash')
                            response = model.generate_content([
                                prompt,
                                {
                                    "mime_type": "image/png",
                                    "data": image_b64
                                }
                            ])
                            
                            if response.text and response.text.strip():
                                analyzed_content = response.text.strip()
                                ocr_texts.append(f"[Slide Image {shape_index + 1} - Gemini]\n{analyzed_content}")
                                logger.info(f"Gemini Vision OCR réussi pour slide image {shape_index + 1}")
                                continue
                                
                        except Exception as e:
                            logger.warning(f"Gemini Vision slide failed: {e}")
                        
                        # Méthode 2: Fallback Tesseract
                        try:
                            import pytesseract
                            image = Image.open(io.BytesIO(image_bytes))
                            text = pytesseract.image_to_string(image, lang='fra+eng')
                            
                            if text.strip():
                                ocr_texts.append(f"[Slide Image {shape_index + 1}] {text.strip()}")
                                
                        except Exception as e2:
                            logger.warning(f"Tesseract slide OCR failed: {e2}")
                            continue
                            
                    except Exception as e:
                        logger.warning(f"Erreur OCR image slide {shape_index}: {e}")
                        continue
            
            return "\n\n".join(ocr_texts)
            
        except Exception as e:
            logger.warning(f"Erreur générale OCR slides: {e}")
            return ""
    
    def _should_apply_ocr(self, content: str, file_path: str) -> bool:
        """Déterminer si l'OCR doit être appliqué"""
        # Configuration OCR
        ocr_config = self.config.get('ocr', {})
        
        # OCR désactivé globalement
        if not ocr_config.get('enabled', False):
            return False
        
        # Types de fichiers supportés
        supported_types = ocr_config.get('supported_types', ['.pdf', '.pptx'])
        file_ext = os.path.splitext(file_path.lower())[1]
        
        if file_ext not in supported_types:
            return False
        
        # Seuil de texte minimum pour déclencher OCR
        min_text_threshold = ocr_config.get('min_text_threshold', 100)
        
        return len(content.strip()) < min_text_threshold
    
    def process_embeddings(self, file_contents: Dict[str, str]) -> List[SimplifiedChunk]:
        """Traiter les documents : chunking + embeddings"""
        logger.info("🧠 Traitement des embeddings...")
        
        # Traitement en batch
        processed_chunks = self.embedding_manager.batch_process_documents(file_contents)
        
        logger.info(f"✅ {len(processed_chunks)} chunks traités")
        
        # Afficher les métriques de performance
        self.embedding_manager.print_metrics()
        
        return processed_chunks
    
    def store_in_vector_db(self, processed_chunks: List[SimplifiedChunk]) -> bool:
        """Stocker les chunks dans la base vectorielle"""
        logger.info("📦 Stockage dans la base vectorielle...")
        
        # Convertir en VectorDocument
        vector_docs = []
        for chunk in processed_chunks:
            vector_doc = VectorDocument(
                id=chunk.id,
                content=chunk.metadata.contenu_chunk,  # Utiliser le contenu original
                embedding=chunk.embedding,
                metadata={
                    "source_file": chunk.metadata.nom_document,
                    "document_type": "simplified",  # Type simplifié
                    "section_title": chunk.enriched_content.theme_principal,
                    "business_context": chunk.enriched_content.commentaire_guidage,
                    "themes": [chunk.enriched_content.theme_principal],
                    "confidence_score": 1.0,  # Score par défaut
                    "chunk_index": chunk.metadata.chunk_index,
                    "token_count": chunk.metadata.taille_chunk,
                    "file_hash": chunk.id.split("_")[0],  # Extraire du ID
                    "created_at": chunk.metadata.date_creation,
                    "dossier_racine": chunk.metadata.dossier_racine,
                    "contenu_enrichi": chunk.enriched_content.contenu_enrichi
                }
            )
            vector_docs.append(vector_doc)
        
        # Stocker en batch
        batch_size = self.config['embeddings']['batch_size']
        total_stored = 0
        
        with tqdm(total=len(vector_docs), desc="Stockage vectoriel") as pbar:
            for i in range(0, len(vector_docs), batch_size):
                batch = vector_docs[i:i + batch_size]
                
                if self.vector_manager.add_documents(batch):
                    total_stored += len(batch)
                    logger.debug(f"Batch {i//batch_size + 1} stocké: {len(batch)} documents")
                else:
                    logger.error(f"Erreur stockage batch {i//batch_size + 1}")
                
                pbar.update(len(batch))
        
        logger.info(f"✅ {total_stored} documents stockés dans la base vectorielle")
        return total_stored == len(vector_docs)
    
    def test_search(self, query: str = "analyse financière") -> List[VectorDocument]:
        """Tester la recherche vectorielle"""
        logger.info(f"🔍 Test de recherche: '{query}'")
        
        # Générer l'embedding de la requête
        query_embedding = self.embedding_manager.get_embedding(query)
        
        # Rechercher
        results = self.vector_manager.search_similar(
            query_embedding=query_embedding,
            top_k=5
        )
        
        logger.info(f"✅ {len(results)} résultats trouvés")
        
        for i, result in enumerate(results, 1):
            logger.info(f"Résultat {i}:")
            logger.info(f"  Score: {result.score:.3f}")
            logger.info(f"  Source: {result.metadata.get('source_file', 'Unknown')}")
            logger.info(f"  Section: {result.metadata.get('section_title', 'Unknown')}")
            logger.info(f"  Contenu: {result.content[:150]}...")
            logger.info("-" * 50)
        
        return results
    
    def run_full_pipeline(self) -> bool:
        """Exécuter le pipeline complet"""
        logger.info("🚀 Démarrage du pipeline RAG complet")
        
        try:
            # 1. Découvrir les documents
            files_info = self.discover_documents()
            if not files_info:
                logger.error("❌ Aucun document trouvé")
                return False
            
            # 2. Télécharger et extraire le contenu
            file_contents = self.download_documents(files_info)
            if not file_contents:
                logger.error("❌ Aucun contenu extrait")
                return False
            
            # 3. Traitement des embeddings
            processed_chunks = self.process_embeddings(file_contents)
            if not processed_chunks:
                logger.error("❌ Aucun chunk traité")
                return False
            
            # 4. Stockage vectoriel
            success = self.store_in_vector_db(processed_chunks)
            if not success:
                logger.error("❌ Erreur stockage vectoriel")
                return False
            
            # 5. Test de recherche
            self.test_search("stratégie entreprise")
            self.test_search("chiffres financiers")
            
            logger.info("🎉 Pipeline RAG terminé avec succès!")
            return True
            
        except Exception as e:
            logger.error(f"❌ Erreur pipeline: {e}")
            return False
    
    def get_stats(self) -> Dict[str, Any]:
        """Obtenir les statistiques du pipeline"""
        metrics = self.embedding_manager.get_metrics()
        return {
            "cache_files": len(list(Path(self.config['google_drive']['cache_directory']).glob("*"))),
            "embedding_cache_size": metrics["cache_sizes"]["embedding_cache"],
            "enrichment_cache_size": metrics["cache_sizes"]["enrichment_cache"],
            "embeddings_generated": metrics["embeddings_generated"],
            "enrichments_generated": metrics["enrichments_generated"],
            "vector_store_type": self.config['vector_store']['type'],
            "index_name": self.config['vector_store']['index_name']
        }

def main():
    """Point d'entrée principal"""
    print("🚀 Pipeline RAG EODEN - Démarrage")
    
    # Vérifier les prérequis
    required_files = [
        "config/settings.yaml",
        "config/credentials.json",
        ".env"
    ]
    
    for file_path in required_files:
        if not Path(file_path).exists():
            print(f"❌ Fichier manquant: {file_path}")
            print("Assurez-vous d'avoir configuré tous les prérequis")
            return
    
    # Créer et exécuter le pipeline
    try:
        pipeline = RAGPipeline()
        
        # Afficher les statistiques initiales
        stats = pipeline.get_stats()
        print(f"📊 Statistiques initiales: {stats}")
        
        # Exécuter le pipeline
        success = pipeline.run_full_pipeline()
        
        if success:
            # Afficher les statistiques finales
            final_stats = pipeline.get_stats()
            print(f"📊 Statistiques finales: {final_stats}")
            print("✅ Pipeline terminé avec succès!")
        else:
            print("❌ Pipeline échoué")
            
    except KeyboardInterrupt:
        print("\n⚠️ Pipeline interrompu par l'utilisateur")
    except Exception as e:
        print(f"❌ Erreur critique: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()